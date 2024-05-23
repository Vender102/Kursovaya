import os
import re
import time

import torch
import gc
from pptx import Presentation
from PyPDF2 import PdfReader
from docx import Document

from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings.huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_community.llms import HuggingFacePipeline
from langchain.prompts import PromptTemplate
from langchain.schema.runnable import RunnablePassthrough
from langchain_community.llms.huggingface_pipeline import HuggingFacePipeline
from langchain.chains import LLMChain

from transformers import (
    AutoModelForCausalLM,
    AutoTokenizer,
    BitsAndBytesConfig,
    pipeline
)

# Activate 4-bit precision base model loading
use_4bit = True
bnb_4bit_compute_dtype = "float16"
bnb_4bit_quant_type = "nf4"
use_nested_quant = False
compute_dtype = getattr(torch, bnb_4bit_compute_dtype)



if compute_dtype == torch.float16 and use_4bit:
    major, _ = torch.cuda.get_device_capability()
    if major >= 8:
        print("=" * 80)
        print("Your GPU supports bfloat16: accelerate training with bf16=True")
        print("=" * 80)

def chunker(research_paper):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1024,
        chunk_overlap=200,
        length_function=len,
    )
    chunked_documents = text_splitter.split_text(research_paper)
    print(len(chunked_documents))
    db = FAISS.from_texts(chunked_documents, HuggingFaceEmbeddings(model_name='sentence-transformers/all-mpnet-base-v2'))
    retriever = db.as_retriever()
    return retriever


def define_model(model, tokenizer, retriever):
    # return tokenizer, model
    text_generation_pipeline = pipeline(
        model=model,
        tokenizer=tokenizer,
        task="text-generation",
        temperature=0.1,
        repetition_penalty=1.1,
        return_full_text=True,
        max_new_tokens=100,

    )

    mistral_llm = HuggingFacePipeline(pipeline=text_generation_pipeline)

    prompt_template = """
        ### [INST] Instruction: The answer the qwestion is based only on knowledge in this context. Give the answer in Russian. Knowing the question, answer only "0" if not or "1" if yes. Here's some context to help:

        {context}

        ### QUESTION:
        {question} [/INST]
        """

    prompt = PromptTemplate(
        input_variables=["context", "question"],
        template=prompt_template,
    )

    llm_chain = LLMChain(llm=mistral_llm, prompt=prompt)

    rag_chain = ({"context": retriever, "question": RunnablePassthrough()} | llm_chain)

    return rag_chain



def extract_text_from_pdf(pdf_file_path):
    try:
        reader = PdfReader(pdf_file_path)
        pages = [page.extract_text() for page in reader.pages]
        document = ''.join(pages)
        ans = []
        for char in document:
            if char.isalnum() or char == ' ' or char == '.' or char == '\n':
                ans.append(char)
        useful_text = ''.join(ans).lower()
        return useful_text
    except Exception as e:
        print(f"Error: {e}")
        return ""

def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    extracted_text = ""
    for paragraph in doc.paragraphs:
        if any(run.text.strip() for run in paragraph.runs):
            extracted_text += paragraph.text + "\n"
    return extracted_text

def extract_text_from_pptx(ppt_file):
    presentation = Presentation(ppt_file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
            elif shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + "\n"
    text = re.sub(r'\s{2,}', ' ', text)
    return text.strip()

def extract_text_from_txt(txt_file_path):
    try:
        with open(txt_file_path, 'r', encoding='utf-8') as file:
            text = file.read()
        return text.strip()
    except Exception as e:
        print(f"Error: {e}")
        return ""


def file_extension(file_list):
    all_data = []
    for file in file_list:
        if file.endswith('.pptx'):
            print('ppt or pptx file')
            data = extract_text_from_pptx(file)
        elif file.endswith('.pdf'):
            print('pdf file')
            data = extract_text_from_pdf(file)
        elif file.endswith('.docx'):
            print('doc or docx file')
            data = extract_text_from_docx(file)
        elif file.endswith('.txt'):
            print('txt')
            data = extract_text_from_txt(file)
        else:
            data = "Unsupported file format"
        all_data.append((file, data))
    return all_data

def write_to_file(filename, data):
    with open(filename, 'a') as f:  # Append mode
        f.write(" ".join(data))

# Get list of files

def sanitize_filename(filename):
    # Replace invalid characters with underscores
    return "".join(c if c.isalnum() or c in "._-" else "_" for c in filename)

def run_model(data_dir, data_out, question, batch_size=1):
    # Load the model and tokenizer only when this function is called
    model_name = 'mistralai/Mistral-7B-Instruct-v0.1'
    use_4bit = True  # Set to True or False depending on your configuration
    bnb_4bit_quant_type = 'nf4'  # Example quantization type
    compute_dtype = torch.float16  # Example compute dtype
    use_nested_quant = True  # Example nested quantization

    bnb_config = BitsAndBytesConfig(
        load_in_4bit=use_4bit,
        bnb_4bit_quant_type=bnb_4bit_quant_type,
        bnb_4bit_compute_dtype=compute_dtype,
        bnb_4bit_use_double_quant=use_nested_quant,
    )

    model = AutoModelForCausalLM.from_pretrained(
        model_name,
        quantization_config=bnb_config,
    )
    tokenizer = AutoTokenizer.from_pretrained(model_name, trust_remote_code=True)
    tokenizer.pad_token = tokenizer.eos_token
    tokenizer.padding_side = "right"

    sanitized_question = sanitize_filename(question)
    output_file_name = os.path.join(data_out, f'question_{sanitized_question}.txt')

    file_list = [os.path.join(data_dir, file) for file in os.listdir(data_dir)]
    for i in range(0, len(file_list), batch_size):
        batch_files = file_list[i:i + batch_size]

        # Extract text from files in the current batch
        research_papers = file_extension(batch_files)
        if not research_papers:
            print("No valid document found in the directory.")
            os._exit(0)

        # Prepare retrievers for the current batch
        retriever_list = []
        for file_name, research_paper in research_papers:
            retriever = chunker(research_paper)
            retriever_list.append((file_name, retriever))

        all_answers = []
        for file_name, retriever in retriever_list:
            rag_chain = define_model(model, tokenizer, retriever)
            result = rag_chain.invoke(question)
            answer = result["text"]
            index = answer.find("[/INST]")
            if index != -1:
                answer = answer[index + len("[/INST]"):].strip()
            all_answers.append(answer)

        # Write all answers to the single output file
        write_to_file(output_file_name, all_answers)
        print(f"Answers for batch saved to {output_file_name}")

        # Clear cache
        del rag_chain, research_papers, retriever_list, all_answers, retriever
        gc.collect()
        torch.cuda.empty_cache()
        print("Cleared GPU cache")

    model = None
    tokenizer = None
    del bnb_config
    gc.collect()
    torch.cuda.empty_cache()



# Ensure no models are loaded at the start
if __name__ == "__main__":
    # Placeholder main to prevent immediate execution
    print("This module provides functions for processing audio files and running models. Do not execute directly.")


